"""
TrustLens — Hackathon Pitch Deck Generator
Generates a premium, dark-themed widescreen presentation.
"""

import os
import sys

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
except ImportError:
    print("Error: python-pptx not installed. Run: pip install python-pptx")
    sys.exit(1)


# ─────────────────────────────────────────────────────────────────
# DESIGN TOKENS
# ─────────────────────────────────────────────────────────────────
BG_DARK      = RGBColor(10, 10, 18)       # Near-black canvas
BG_CARD      = RGBColor(22, 27, 45)       # Card surfaces
BG_CARD_ALT  = RGBColor(28, 34, 55)       # Slightly lighter card
TEAL         = RGBColor(0, 224, 196)      # Primary accent
PURPLE       = RGBColor(168, 130, 255)    # Secondary accent
AMBER        = RGBColor(255, 183, 77)     # Warning / highlight
RED_SOFT     = RGBColor(255, 99, 99)      # Danger / reject
GREEN_SOFT   = RGBColor(74, 222, 128)     # Approve / success
WHITE        = RGBColor(245, 245, 250)
GRAY         = RGBColor(140, 150, 175)
GRAY_DIM     = RGBColor(90, 100, 120)
SLIDE_W      = Inches(13.333)
SLIDE_H      = Inches(7.5)


def _bg(slide):
    """Apply dark background to a slide."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG_DARK


def _rect(slide, left, top, w, h, color, alpha=None):
    """Add a solid colour rectangle (decorative)."""
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    return sh


def _rounded_card(slide, left, top, w, h, border_color=None, bg=None):
    """Add a rounded-rectangle card with optional border."""
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = bg or BG_CARD
    if border_color:
        sh.line.color.rgb = border_color
        sh.line.width = Pt(1.5)
    else:
        sh.line.fill.background()
    return sh


def _text(slide, left, top, w, h):
    """Add a text box and return its text_frame."""
    tb = slide.shapes.add_textbox(left, top, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    return tf


def _para(tf, text, size=12, bold=False, color=WHITE, align=None,
          spacing_after=6, line_spacing=None, italic=False, first=False):
    """Add a paragraph to a text_frame."""
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.text = text
    p.font.name = "Arial"
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.italic = italic
    p.font.color.rgb = color
    p.space_after = Pt(spacing_after)
    if align:
        p.alignment = align
    if line_spacing:
        p.line_spacing = line_spacing
    return p


def _circle(slide, left, top, size, color):
    """Add a decorative circle."""
    sh = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    return sh


def _slide_number(slide, num, total=8):
    """Add a subtle slide number at the bottom-right."""
    tf = _text(slide, Inches(11.8), Inches(7.0), Inches(1.0), Inches(0.35))
    _para(tf, f"{num} / {total}", size=9, color=GRAY_DIM,
          align=PP_ALIGN.RIGHT, first=True)


def _header_bar(slide, label, title, num, total=8):
    """Standard slide header with category label and title."""
    # Thin accent line at top
    _rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), TEAL)

    # Category label
    tf = _text(slide, Inches(0.9), Inches(0.35), Inches(10), Inches(0.3))
    _para(tf, label.upper(), size=9, bold=True, color=TEAL,
          spacing_after=0, first=True)

    # Title
    tf2 = _text(slide, Inches(0.9), Inches(0.6), Inches(10), Inches(0.7))
    _para(tf2, title, size=30, bold=True, color=WHITE,
          spacing_after=0, first=True)

    # Team name top-right
    tf3 = _text(slide, Inches(10.5), Inches(0.35), Inches(2.5), Inches(0.3))
    _para(tf3, "TEAM ALONE", size=9, bold=True, color=PURPLE,
          align=PP_ALIGN.RIGHT, spacing_after=0, first=True)

    _slide_number(slide, num, total)


# ═════════════════════════════════════════════════════════════════
# BUILD DECK
# ═════════════════════════════════════════════════════════════════
def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]

    # ─────────────────────────────────────────────────────────────
    # SLIDE 1 — COVER
    # ─────────────────────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    _bg(s)

    # Decorative shapes
    _rect(s, Inches(0), Inches(0), Inches(0.12), SLIDE_H, TEAL)
    _circle(s, Inches(10.2), Inches(0.6), Inches(1.6), BG_CARD)
    _circle(s, Inches(10.6), Inches(1.0), Inches(0.8), BG_CARD_ALT)

    # Gradient-like decorative bar behind title
    _rect(s, Inches(0.12), Inches(2.0), Inches(0.5), Inches(3.5), BG_CARD)

    # Team name — top right corner
    tf_team = _text(s, Inches(9.5), Inches(0.5), Inches(3.5), Inches(0.5))
    _para(tf_team, "TEAM ALONE", size=13, bold=True, color=PURPLE,
          align=PP_ALIGN.RIGHT, first=True)

    # Main content
    tf = _text(s, Inches(1.2), Inches(1.6), Inches(8.5), Inches(5.0))

    # Tagline
    _para(tf, "AI-POWERED REVIEW VERIFICATION", size=11, bold=True,
          color=TEAL, spacing_after=10, first=True)

    # Project name
    p_name = tf.add_paragraph()
    p_name.text = "TrustLens"
    p_name.font.name = "Arial"
    p_name.font.size = Pt(72)
    p_name.font.bold = True
    p_name.font.color.rgb = WHITE
    p_name.space_after = Pt(6)

    # Emoji shield line
    _para(tf, "🛡️  Detect Fake Reviews  ·  Verify Purchases  ·  Analyze Images  ·  Score Trust",
          size=16, color=GRAY, spacing_after=20)

    # Tech line
    _para(tf, "Built with  LangGraph  ·  LangChain  ·  Groq Llama 3.3 70B  ·  Llama 4 Scout Vision  ·  FastAPI  ·  React 19",
          size=12, color=GRAY_DIM, spacing_after=30)

    # Separator
    _rect(s, Inches(1.2), Inches(5.4), Inches(4.5), Inches(0.03), TEAL)

    # Stats row
    stats = [("7", "AI Agents"), ("3", "Custom Tools"), ("100%", "Transparent Score")]
    for i, (val, lbl) in enumerate(stats):
        x = Inches(1.2 + i * 2.6)
        tf_s = _text(s, x, Inches(5.6), Inches(2.2), Inches(1.2))
        _para(tf_s, val, size=32, bold=True, color=TEAL, spacing_after=2, first=True)
        _para(tf_s, lbl, size=11, color=GRAY, spacing_after=0)

    _slide_number(s, 1)

    # ─────────────────────────────────────────────────────────────
    # SLIDE 2 — PROBLEM STATEMENT
    # ─────────────────────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    _bg(s)
    _header_bar(s, "The Problem", "Fake Reviews Are Destroying Digital Trust", 2)

    # Subtitle
    tf_sub = _text(s, Inches(0.9), Inches(1.35), Inches(11), Inches(0.5))
    _para(tf_sub, "Every day, millions of consumers are misled by fabricated reviews — costing businesses billions and eroding marketplace integrity.",
          size=13, color=GRAY, first=True, line_spacing=1.3)

    # Three problem cards
    problems = [
        ("💰", "Review Fraud\n& Astroturfing",
         "Competitors deploy bot armies and paid reviewers to manipulate star ratings. A single bad actor can submit hundreds of fake reviews overnight.",
         TEAL),
        ("📷", "Visual\nMisrepresentation",
         "Fake reviewers upload stock photos, AI-generated images, or screenshots from other platforms to appear as genuine verified buyers.",
         PURPLE),
        ("⏱️", "Unscalable\nManual Auditing",
         "Cross-referencing receipts, checking text consistency, and validating image authenticity manually is slow, expensive, and error-prone at scale.",
         AMBER),
    ]

    for i, (icon, title, desc, color) in enumerate(problems):
        left = Inches(0.9 + i * 4.0)
        card = _rounded_card(s, left, Inches(2.2), Inches(3.7), Inches(4.6), border_color=color)
        ctf = card.text_frame
        ctf.word_wrap = True
        ctf.margin_left = Inches(0.3)
        ctf.margin_right = Inches(0.3)
        ctf.margin_top = Inches(0.35)

        # Icon
        _para(ctf, icon, size=36, spacing_after=8, first=True)
        # Title
        _para(ctf, title, size=18, bold=True, color=color, spacing_after=12)
        # Desc
        _para(ctf, desc, size=12, color=GRAY, line_spacing=1.35, spacing_after=14)

    # Bottom stat bar
    stat_bar = _rounded_card(s, Inches(0.9), Inches(7.0), Inches(11.5), Inches(0.35), bg=BG_CARD_ALT)
    stat_tf = stat_bar.text_frame
    stat_tf.word_wrap = True
    stat_tf.margin_left = Inches(0.3)
    stat_tf.margin_top = Inches(0.05)
    _para(stat_tf, "🔴  $152B lost annually to fake reviews   ·   42% of online reviews are suspect   ·   82% of consumers trust reviews like personal recommendations",
          size=10, color=GRAY, first=True)

    # ─────────────────────────────────────────────────────────────
    # SLIDE 3 — INTRODUCING TRUSTLENS
    # ─────────────────────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    _bg(s)
    _header_bar(s, "The Solution", "Introducing TrustLens", 3)

    # Left description block
    tf_l = _text(s, Inches(0.9), Inches(1.5), Inches(5.2), Inches(5.5))
    _para(tf_l, "Collective AI Auditing", size=22, bold=True, color=TEAL,
          spacing_after=12, first=True)
    _para(tf_l,
          "TrustLens is a full-stack AI review verification engine that treats every review as a case to be audited by a team of specialized AI agents.\n\n"
          "Instead of relying on simple heuristics or keyword filters, TrustLens orchestrates 7 intelligent agents inside a LangGraph state machine. "
          "Each agent examines one dimension of trust — from invoice verification to image forensics — before a final Trust Score agent synthesizes "
          "all findings into a transparent 0-100 verdict.",
          size=13, color=GRAY, line_spacing=1.4, spacing_after=20)

    _para(tf_l, "✦  Every score is database-backed and explainable",
          size=12, color=WHITE, spacing_after=6)
    _para(tf_l, "✦  Real-time image analysis via Llama 4 Scout Vision",
          size=12, color=WHITE, spacing_after=6)
    _para(tf_l, "✦  Full audit trail with per-agent breakdown",
          size=12, color=WHITE, spacing_after=6)
    _para(tf_l, "✦  Sub-10 second end-to-end verification",
          size=12, color=WHITE, spacing_after=6)

    # Right — 4 pillar cards (stacked)
    pillars = [
        ("🗄️  Database Verification", "Cross-references bill IDs against live SQLite order records with customer, item, and amount data.", TEAL),
        ("🧠  LangGraph Orchestration", "7 agents wired in a directed state graph with conditional routing, state passing, and parallel checks.", PURPLE),
        ("📷  Vision AI Analysis", "Llama 4 Scout analyzes uploaded photos — detects stock images, AI generation, and context mismatches.", TEAL),
        ("📊  Transparent Trust Score", "Weighted 0-100 score with full per-category breakdown. Clear Approve / Review / Reject verdict.", PURPLE),
    ]

    for i, (title, desc, color) in enumerate(pillars):
        top = Inches(1.5 + i * 1.4)
        card = _rounded_card(s, Inches(6.5), top, Inches(6.0), Inches(1.2), border_color=color)
        ctf = card.text_frame
        ctf.word_wrap = True
        ctf.margin_left = Inches(0.25)
        ctf.margin_right = Inches(0.25)
        ctf.margin_top = Inches(0.15)

        _para(ctf, title, size=13, bold=True, color=color, spacing_after=4, first=True)
        _para(ctf, desc, size=11, color=GRAY, line_spacing=1.2, spacing_after=0)

    # ─────────────────────────────────────────────────────────────
    # SLIDE 4 — END-TO-END PIPELINE
    # ─────────────────────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    _bg(s)
    _header_bar(s, "Architecture", "End-to-End Verification Pipeline", 4)

    # Subtitle
    tf_sub = _text(s, Inches(0.9), Inches(1.35), Inches(11), Inches(0.45))
    _para(tf_sub, "From review submission to Trust Score verdict — every step is orchestrated by LangGraph with full state management and tracing via LangSmith.",
          size=12, color=GRAY, first=True)

    # Pipeline flow — horizontal connected boxes
    pipeline_steps = [
        ("📥", "Review\nSubmitted", "React Frontend\nCaptures review + images", TEAL),
        ("⚡", "FastAPI\nBackend", "REST endpoint\nFile upload & routing", PURPLE),
        ("🔄", "LangGraph\nEngine", "State machine\nOrchestrates 7 agents", TEAL),
        ("🗄️", "SQLite\nDatabase", "Order verification\nReview persistence", PURPLE),
        ("⭐", "Trust\nScore", "0–100 verdict\nApprove / Review / Reject", AMBER),
    ]

    for i, (icon, title, desc, color) in enumerate(pipeline_steps):
        left = Inches(0.5 + i * 2.5)
        card = _rounded_card(s, left, Inches(2.0), Inches(2.2), Inches(2.3), border_color=color)
        ctf = card.text_frame
        ctf.word_wrap = True
        ctf.margin_left = Inches(0.15)
        ctf.margin_right = Inches(0.15)
        ctf.margin_top = Inches(0.2)

        _para(ctf, icon, size=28, spacing_after=4, align=PP_ALIGN.CENTER, first=True)
        _para(ctf, title, size=14, bold=True, color=color, spacing_after=6, align=PP_ALIGN.CENTER)
        _para(ctf, desc, size=10, color=GRAY, align=PP_ALIGN.CENTER, line_spacing=1.2)

        # Arrow connector (except last)
        if i < len(pipeline_steps) - 1:
            arr_tf = _text(s, left + Inches(2.2), Inches(2.8), Inches(0.3), Inches(0.5))
            _para(arr_tf, "→", size=22, bold=True, color=GRAY_DIM, first=True)

    # Bottom detail cards — left and right
    # Left: Agent flow
    lc = _rounded_card(s, Inches(0.9), Inches(4.6), Inches(5.6), Inches(2.5), border_color=TEAL)
    lctf = lc.text_frame
    lctf.word_wrap = True
    lctf.margin_left = Inches(0.3)
    lctf.margin_top = Inches(0.25)
    _para(lctf, "Agent Execution Order", size=15, bold=True, color=TEAL, spacing_after=10, first=True)
    _para(lctf, "1. Intake Agent → Analyzes review, decides which checks to run", size=11, color=WHITE, spacing_after=5)
    _para(lctf, "2. Purchase Agent → Verifies bill ID in database", size=11, color=WHITE, spacing_after=5)
    _para(lctf, "3. Consistency Agent → Compares text vs purchased items", size=11, color=WHITE, spacing_after=5)
    _para(lctf, "4. Duplicate Agent → Checks for copy-paste plagiarism", size=11, color=WHITE, spacing_after=5)
    _para(lctf, "5. Behavior Agent → Profiles submission patterns", size=11, color=WHITE, spacing_after=5)
    _para(lctf, "6. Vision Agent → Analyzes uploaded images with Llama 4 Scout", size=11, color=WHITE, spacing_after=5)
    _para(lctf, "7. Trust Score Agent → Aggregates all signals → Final verdict", size=11, color=AMBER, spacing_after=0)

    # Right: Tech specs
    rc = _rounded_card(s, Inches(6.8), Inches(4.6), Inches(5.7), Inches(2.5), border_color=PURPLE)
    rctf = rc.text_frame
    rctf.word_wrap = True
    rctf.margin_left = Inches(0.3)
    rctf.margin_top = Inches(0.25)
    _para(rctf, "Infrastructure Details", size=15, bold=True, color=PURPLE, spacing_after=10, first=True)
    _para(rctf, "•  LangGraph StateGraph with 8 nodes (7 agents + DB save)", size=11, color=WHITE, spacing_after=5)
    _para(rctf, "•  3 Groq API keys distributed across agents for rate-limit safety", size=11, color=WHITE, spacing_after=5)
    _para(rctf, "•  Conditional edge routing — agents can be skipped dynamically", size=11, color=WHITE, spacing_after=5)
    _para(rctf, "•  Full LangSmith tracing for every pipeline execution", size=11, color=WHITE, spacing_after=5)
    _para(rctf, "•  SQLite auto-seeds demo orders/businesses on startup", size=11, color=WHITE, spacing_after=5)
    _para(rctf, "•  Uploaded images stored locally, analyzed via base64 encoding", size=11, color=WHITE, spacing_after=0)

    # ─────────────────────────────────────────────────────────────
    # SLIDE 5 — AI AGENTS DEEP DIVE
    # ─────────────────────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    _bg(s)
    _header_bar(s, "AI Agents", "7 Specialized LLM-Powered Agents", 5)

    tf_sub = _text(s, Inches(0.9), Inches(1.35), Inches(11), Inches(0.35))
    _para(tf_sub, "Each agent is a specialized LLM node powered by Groq Llama 3.3 70B. The Vision agent uses Llama 4 Scout for multimodal image analysis.",
          size=12, color=GRAY, first=True)

    agents = [
        ("📥", "01", "Intake Agent",
         "Analyzes the submitted review text and context. Decides which verification agents to activate based on content, bill ID presence, and media availability.",
         "Key 1", TEAL),
        ("💳", "02", "Purchase Agent",
         "Queries the SQLite database to verify if the bill ID corresponds to a real order. Validates customer name, items purchased, amount, and order date.",
         "Key 1", PURPLE),
        ("🔍", "03", "Consistency Agent",
         "Compares the review text against verified purchase details. Flags mismatches like reviewing coffee when they bought electronics.",
         "Key 2", TEAL),
        ("📋", "04", "Duplicate Agent",
         "Scans the review database for copy-paste plagiarism. Detects spun content, identical phrases, and reviews reused across different orders.",
         "Key 2", PURPLE),
        ("👤", "05", "Behavior Agent",
         "Profiles the reviewer's submission history — frequency, timing patterns, and review count. Flags bot-like velocity or review farm indicators.",
         "Key 3", TEAL),
        ("📷", "06", "Vision Agent",
         "Uses Llama 4 Scout vision model to analyze uploaded images. Describes photo contents, checks if they match the review, and detects stock/AI images.",
         "Key 3", PURPLE),
        ("⭐", "07", "Trust Score Agent",
         "Aggregates all 6 agent signals into a weighted 0-100 trust score. Issues High/Medium/Low verdict with full point breakdown and action recommendation.",
         "Key 3", AMBER),
    ]

    for i, (icon, num, name, desc, key, color) in enumerate(agents):
        row = i // 4
        col = i % 4

        left = Inches(0.5 + col * 3.1)
        top = Inches(1.9 + row * 2.7)
        w = Inches(2.9)
        h = Inches(2.5)

        card = _rounded_card(s, left, top, w, h, border_color=color)
        ctf = card.text_frame
        ctf.word_wrap = True
        ctf.margin_left = Inches(0.2)
        ctf.margin_right = Inches(0.15)
        ctf.margin_top = Inches(0.15)

        # Header line: icon + number
        _para(ctf, f"{icon}  {num}", size=12, color=GRAY_DIM, spacing_after=2, first=True)
        # Name
        _para(ctf, name, size=14, bold=True, color=color, spacing_after=6)
        # Description
        _para(ctf, desc, size=10, color=GRAY, line_spacing=1.25, spacing_after=6)
        # API key badge
        _para(ctf, f"Groq {key}", size=8, color=GRAY_DIM, italic=True, spacing_after=0)

    # ─────────────────────────────────────────────────────────────
    # SLIDE 6 — TRUST SCORE CALCULATION
    # ─────────────────────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    _bg(s)
    _header_bar(s, "Scoring", "How the Trust Score is Calculated", 6)

    # LEFT — Score breakdown with visual bars
    lc = _rounded_card(s, Inches(0.9), Inches(1.5), Inches(5.6), Inches(5.5), border_color=TEAL)
    lctf = lc.text_frame
    lctf.word_wrap = True
    lctf.margin_left = Inches(0.35)
    lctf.margin_right = Inches(0.35)
    lctf.margin_top = Inches(0.3)

    _para(lctf, "Score Composition  (0 – 100 points)", size=16, bold=True,
          color=TEAL, spacing_after=14, first=True)

    # With media
    _para(lctf, "▎ WITH MEDIA UPLOADED", size=10, bold=True, color=AMBER, spacing_after=8)

    weights_media = [
        ("Purchase Verification", "25 pts", "Bill ID verified in real database", TEAL),
        ("Experience Consistency", "25 pts", "Review text matches purchased items", PURPLE),
        ("Duplicate Detection", "15 pts", "No plagiarism or copy-paste detected", TEAL),
        ("Reviewer Behavior", "15 pts", "No bot patterns or spam velocity", PURPLE),
        ("Media Authenticity", "20 pts", "Photos match review + not stock/AI", AMBER),
    ]
    for name, pts, desc, color in weights_media:
        _para(lctf, f"{'━' * 3}  {name}  ·  {pts}", size=11, bold=True, color=color, spacing_after=1)
        _para(lctf, f"      {desc}", size=10, color=GRAY_DIM, spacing_after=8)

    _para(lctf, "▎ WITHOUT MEDIA  →  Points redistribute:", size=10, bold=True, color=AMBER, spacing_after=4)
    _para(lctf, "Purchase 30 pts  ·  Consistency 30 pts  ·  Duplicate 20 pts  ·  Behavior 20 pts",
          size=10, color=GRAY, spacing_after=0)

    # RIGHT — Verdict tiers
    rc = _rounded_card(s, Inches(6.8), Inches(1.5), Inches(5.7), Inches(5.5), border_color=PURPLE)
    rctf = rc.text_frame
    rctf.word_wrap = True
    rctf.margin_left = Inches(0.35)
    rctf.margin_right = Inches(0.35)
    rctf.margin_top = Inches(0.3)

    _para(rctf, "Verdict Tiers & Actions", size=16, bold=True, color=PURPLE,
          spacing_after=16, first=True)

    # HIGH
    _para(rctf, "▎ 75 – 100  ·  HIGH TRUST", size=14, bold=True, color=GREEN_SOFT, spacing_after=4)
    _para(rctf, "✅  Action: APPROVE", size=12, bold=True, color=GREEN_SOFT, spacing_after=2)
    _para(rctf, "Order verified  ·  Consistent experience  ·  No duplicates or bot patterns  ·  Media matches (if uploaded)", size=10, color=GRAY, spacing_after=16)

    # MEDIUM
    _para(rctf, "▎ 45 – 74  ·  MEDIUM TRUST", size=14, bold=True, color=AMBER, spacing_after=4)
    _para(rctf, "⚠️  Action: MANUAL REVIEW", size=12, bold=True, color=AMBER, spacing_after=2)
    _para(rctf, "Some verification checkpoints are weak or incomplete. One or more agents flagged minor concerns.", size=10, color=GRAY, spacing_after=16)

    # LOW
    _para(rctf, "▎ 0 – 44  ·  LOW TRUST", size=14, bold=True, color=RED_SOFT, spacing_after=4)
    _para(rctf, "🚫  Action: REJECT", size=12, bold=True, color=RED_SOFT, spacing_after=2)
    _para(rctf, "Multiple verification checks failed. Highly likely to be a fake, fabricated, or spam review.", size=10, color=GRAY, spacing_after=16)

    # Example
    _para(rctf, "Example Output:", size=11, bold=True, color=WHITE, spacing_after=4)
    _para(rctf, '{ "score": 87, "level": "High", "action": "approve",\n  "breakdown": { purchase: 25, consistency: 22,\n    duplicate: 15, behavior: 13, media: 12 } }',
          size=9, color=GRAY_DIM, spacing_after=0, line_spacing=1.2)

    # ─────────────────────────────────────────────────────────────
    # SLIDE 7 — LINKS & DEMO
    # ─────────────────────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    _bg(s)
    _header_bar(s, "Live Demo", "Try TrustLens — Links & Resources", 7)

    # Center card
    center_card = _rounded_card(s, Inches(2.5), Inches(1.8), Inches(8.3), Inches(4.8), border_color=TEAL)
    cctf = center_card.text_frame
    cctf.word_wrap = True
    cctf.margin_left = Inches(0.6)
    cctf.margin_right = Inches(0.6)
    cctf.margin_top = Inches(0.5)

    _para(cctf, "🔗  Project Links", size=24, bold=True, color=WHITE,
          spacing_after=24, first=True, align=PP_ALIGN.CENTER)

    links = [
        ("GitHub Repository", "github.com/Abuzaid-01/TrustLens", TEAL),
        ("Live Backend (Render)", "[ Add deployed URL here ]", PURPLE),
        ("Live Frontend", "[ Add deployed URL here ]", TEAL),
        ("Demo Video", "[ Add video link here ]", PURPLE),
    ]
    for title, url, color in links:
        _para(cctf, title, size=14, bold=True, color=color, spacing_after=2, align=PP_ALIGN.CENTER)
        _para(cctf, url, size=12, color=GRAY, spacing_after=20, align=PP_ALIGN.CENTER)

    # Tech badges at bottom
    tf_badges = _text(s, Inches(1.5), Inches(6.8), Inches(10), Inches(0.4))
    _para(tf_badges, "LangGraph  ·  LangChain  ·  Groq  ·  Llama 3.3 70B  ·  Llama 4 Scout  ·  LangSmith  ·  FastAPI  ·  React 19  ·  SQLite",
          size=10, color=GRAY_DIM, first=True, align=PP_ALIGN.CENTER)

    # ─────────────────────────────────────────────────────────────
    # SLIDE 8 — WHY TRUSTLENS / BENEFITS / CLOSING
    # ─────────────────────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    _bg(s)
    _header_bar(s, "Why TrustLens", "Benefits & Impact", 8)

    tf_sub = _text(s, Inches(0.9), Inches(1.35), Inches(11), Inches(0.4))
    _para(tf_sub, "TrustLens doesn't just detect fake reviews — it builds a transparent audit infrastructure that platforms and consumers can trust.",
          size=13, color=GRAY, first=True)

    benefits = [
        ("🔒", "Restore Platform Trust",
         "Every review score is backed by real database verification and transparent per-category breakdown. Customers buy with confidence.",
         TEAL),
        ("⚡", "95% Cost Reduction",
         "Automates the entire receipt verification, text analysis, image checking, and behavior profiling pipeline. Replaces manual moderation teams.",
         PURPLE),
        ("🧩", "Extensible Architecture",
         "LangGraph's stateful nodes make it trivial to add new agents (sentiment analysis, language detection, cross-platform checks) without rewriting the pipeline.",
         TEAL),
        ("📷", "Multimodal Protection",
         "Shills can't fake proof-of-ownership with stock photos. Llama 4 Scout vision model catches mismatches between uploaded images and review claims.",
         PURPLE),
        ("📊", "Full Transparency",
         "Unlike black-box spam filters, TrustLens provides a complete audit trail — every agent's reasoning, every score component, every flag raised.",
         AMBER),
        ("🌐", "Production Ready",
         "FastAPI backend, React frontend, SQLite database, Groq cloud inference. Full-stack, deployable today, scalable to PostgreSQL + cloud tomorrow.",
         GREEN_SOFT),
    ]

    for i, (icon, title, desc, color) in enumerate(benefits):
        row = i // 3
        col = i % 3
        left = Inches(0.7 + col * 4.0)
        top = Inches(2.0 + row * 2.5)

        card = _rounded_card(s, left, top, Inches(3.8), Inches(2.3), border_color=color)
        ctf = card.text_frame
        ctf.word_wrap = True
        ctf.margin_left = Inches(0.25)
        ctf.margin_right = Inches(0.2)
        ctf.margin_top = Inches(0.2)

        _para(ctf, f"{icon}  {title}", size=14, bold=True, color=color,
              spacing_after=8, first=True)
        _para(ctf, desc, size=11, color=GRAY, line_spacing=1.3, spacing_after=0)

    # Closing tagline
    tf_close = _text(s, Inches(2), Inches(7.0), Inches(9), Inches(0.4))
    _para(tf_close, "TrustLens  —  Making Every Review Count  🛡️",
          size=14, bold=True, color=WHITE, first=True, align=PP_ALIGN.CENTER)

    # ─────────────────────────────────────────────────────────────
    # SAVE
    # ─────────────────────────────────────────────────────────────
    out = "TrustLens_Presentation.pptx"
    prs.save(out)
    print(f"✅ Presentation saved: {os.path.abspath(out)}")


if __name__ == "__main__":
    build()